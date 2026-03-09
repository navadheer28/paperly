import sys
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
import io

def pdf_to_ppt(input_path, output_path):
    doc = fitz.open(input_path)
    prs = Presentation()

    for i in range(len(doc)):
        page = doc[i]
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes('png')

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        width = prs.slide_width
        height = prs.slide_height

        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, 0, 0, width, height)

    doc.close()
    prs.save(output_path)
    print("Converted successfully")

if __name__ == '__main__':
    pdf_to_ppt(sys.argv[1], sys.argv[2])